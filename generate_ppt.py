from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TEAL = RGBColor(0x00, 0x5A, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG = RGBColor(0xF0, 0xF6, 0xF8)
LT_TEAL = RGBColor(0xCC, 0xE5, 0xEB)
ACCENT = RGBColor(0xE8, 0xF4, 0xF8)
ORANGE = RGBColor(0xE8, 0x7C, 0x2A)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)

def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_tb(slide, left, top, width, height, text, sz=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, fname="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = fname
    p.alignment = align
    return tb

def add_bullets(slide, left, top, width, height, items, sz=13, color=DARK, sp=Pt(4)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = sp
    return tb

def add_stat_box(slide, left, top, w, h, number, label, num_color=TEAL):
    box = add_rect(slide, left, top, w, h, ACCENT)
    add_tb(slide, left, top + Inches(0.1), w, Inches(0.55), number, sz=28, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    add_tb(slide, left, top + Inches(0.6), w, Inches(0.4), label, sz=10, color=DARK, align=PP_ALIGN.CENTER)

def make_header(slide, slide_label):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.5), TEAL)
    add_tb(slide, Inches(0.4), Inches(0.02), Inches(8), Inches(0.45),
        "thinkQbator  \u2022  FAST-TRACK BATCH", sz=13, bold=True, color=WHITE)
    add_tb(slide, Inches(8.5), Inches(0.02), Inches(4.5), Inches(0.45),
        "Product Readiness Submission  |  " + slide_label, sz=11, color=LT_TEAL, align=PP_ALIGN.RIGHT)

def make_footer(slide):
    add_rect(slide, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), TEAL)
    add_tb(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4),
        "thinkQbator Fast-Track Batch  |  Local Clothing Store Discovery  |  Confidential",
        sz=10, color=LT_TEAL)

def add_img(slide, path, left, top, width, height):
    if os.path.exists(path):
        try:
            slide.shapes.add_picture(path, left, top, width, height)
            return True
        except:
            return False
    return False

base_dir = r"C:\Users\pramo\Email Automation-Checker\ppt"

# ================================================================
# SLIDE 1 — How your idea works on the field
# ================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
make_header(s1, "Slide 1 of 3")

add_tb(s1, Inches(0.4), Inches(0.6), Inches(12.5), Inches(0.6),
    "How your idea works on the field", sz=28, bold=True, color=TEAL)

# --- LEFT COLUMN ---
lx = Inches(0.4)
cw = Inches(5.8)

# User Journey
add_rect(s1, lx, Inches(1.4), Inches(2.0), Inches(0.35), TEAL)
add_tb(s1, lx + Inches(0.1), Inches(1.42), Inches(1.8), Inches(0.3),
    "User Journey", sz=14, bold=True, color=WHITE)

add_bullets(s1, lx, Inches(1.85), cw, Inches(2.1), [
    "1.  Customer opens app \u2014 location detected via GPS, nearby",
    "    stores shown on map/list view",
    "2.  Customer searches and filters by style, size, category, or price",
    "3.  Customer checks real-time stock and photos for specific items",
    "4.  Customer places order or reserves item for pickup/delivery,",
    "    pays via integrated payment gateway",
    "5.  Store owner gets notification, confirms and fulfills order",
    "    through merchant dashboard",
], sz=12)

# Deployment Details
add_rect(s1, lx, Inches(4.05), Inches(2.5), Inches(0.35), TEAL)
add_tb(s1, lx + Inches(0.1), Inches(4.07), Inches(2.3), Inches(0.3),
    "Deployment Details", sz=14, bold=True, color=WHITE)

add_bullets(s1, lx, Inches(4.5), cw, Inches(1.5), [
    "Mobile app (customer-facing) + Web app + Merchant dashboard",
    "Backend server + Database + Cloud hosting",
    "GPS/Location services for store discovery",
    "Product inventory management system",
    "Payment gateway integration",
    "Notification system (push / in-app)",
    "Search & filter technology",
    "API integrations & basic analytics",
], sz=12)

# --- RIGHT COLUMN ---
rx = Inches(6.8)
rw = Inches(6.1)

# Metrics
add_rect(s1, rx, Inches(1.4), Inches(2.5), Inches(0.35), TEAL)
add_tb(s1, rx + Inches(0.1), Inches(1.42), Inches(2.3), Inches(0.3),
    "Metrics & Impact (Target \u2014 Year 1)", sz=14, bold=True, color=WHITE)

stat_data = [
    ("50+", "Stores\nOnboarded"),
    ("1,000+", "Active\nUsers"),
    ("60%", "Faster\nDiscovery"),
    ("3+", "Localities\nCovered"),
    ("10+", "Pilot\nStores"),
    ("4.5\u2605", "Avg.\nRating"),
]
box_w = Inches(1.85)
box_h = Inches(1.05)
gap = Inches(0.1)
for i, (num, lbl) in enumerate(stat_data):
    row = i // 3
    col = i % 3
    x = rx + col * (box_w + gap)
    y = Inches(1.95) + row * (box_h + gap)
    add_stat_box(s1, x, y, box_w, box_h, num, lbl)

# App mockup screenshot from demo app
add_img(s1, os.path.join(base_dir, "screenshot_home.png"), rx, Inches(4.1), Inches(1.8), Inches(2.6))

add_tb(s1, rx + Inches(1.9), Inches(4.3), Inches(4.2), Inches(2.2),
    "Live interactive prototype built as a mobile-first web app. "
    "Features: GPS-based store discovery, real-time inventory, "
    "cart management, merchant dashboard, and profile management.",
    sz=11, color=DARK)

# Team + Social Impact
add_rect(s1, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.65), ACCENT)
add_tb(s1, Inches(0.5), Inches(6.22), Inches(5), Inches(0.25),
    "Team: Pramod Yadav  |  Mohd Uzair Khan", sz=11, bold=True, color=TEAL)
add_tb(s1, Inches(0.5), Inches(6.48), Inches(12), Inches(0.25),
    "Social Impact: Yes \u2014 Empowers small/local clothing retailers with digital visibility, "
    "helping them compete with large marketplaces and sustain local livelihoods.",
    sz=10, color=DARK)

make_footer(s1)

# ================================================================
# SLIDE 2 — Show the product
# ================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)
make_header(s2, "Slide 2 of 3")

add_tb(s2, Inches(0.4), Inches(0.6), Inches(12.5), Inches(0.6),
    "Show the product", sz=28, bold=True, color=TEAL)
add_tb(s2, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.3),
    "Interactive demo prototype built with HTML/CSS/JS \u2014 open in browser to experience the full app.",
    sz=12, color=GRAY)

# --- Photo (Left) ---
add_rect(s2, Inches(0.4), Inches(1.55), Inches(1.6), Inches(0.30), TEAL)
add_tb(s2, Inches(0.5), Inches(1.56), Inches(1.4), Inches(0.25),
    "Photo \u2014 App Screenshots", sz=12, bold=True, color=WHITE)

photo_box = add_rect(s2, Inches(0.4), Inches(1.9), Inches(6.2), Inches(3.6), LIGHT_BG)
photo_box.line.color.rgb = TEAL
photo_box.line.width = Pt(1.5)

# Home screen + Store detail side by side
add_img(s2, os.path.join(base_dir, "screenshot_home.png"), Inches(0.7), Inches(2.15), Inches(2.0), Inches(3.1))
add_img(s2, os.path.join(base_dir, "screenshot_store.png"), Inches(3.2), Inches(2.15), Inches(2.0), Inches(3.1))

add_tb(s2, Inches(0.4), Inches(5.6), Inches(6.2), Inches(0.3),
    "Home screen with map + store discovery  |  Store detail with product inventory",
    sz=10, color=GRAY, align=PP_ALIGN.CENTER)

# --- Video (Top Right) ---
add_rect(s2, Inches(6.9), Inches(1.55), Inches(1.6), Inches(0.30), TEAL)
add_tb(s2, Inches(7.0), Inches(1.56), Inches(1.4), Inches(0.25),
    "More Screens", sz=12, bold=True, color=WHITE)

more_box = add_rect(s2, Inches(6.9), Inches(1.9), Inches(6.0), Inches(1.7), LIGHT_BG)
more_box.line.color.rgb = TEAL
more_box.line.width = Pt(1.5)

add_img(s2, os.path.join(base_dir, "screenshot_cart.png"), Inches(7.1), Inches(2.05), Inches(1.6), Inches(1.4))
add_img(s2, os.path.join(base_dir, "screenshot_merchant.png"), Inches(8.9), Inches(2.05), Inches(1.6), Inches(1.4))
add_img(s2, os.path.join(base_dir, "screenshot_profile.png"), Inches(10.7), Inches(2.05), Inches(1.6), Inches(1.4))

add_tb(s2, Inches(6.9), Inches(3.7), Inches(6.0), Inches(0.3),
    "Cart  |  Merchant Dashboard  |  Profile",
    sz=10, color=GRAY, align=PP_ALIGN.CENTER)

# Hyperlink
add_rect(s2, Inches(6.9), Inches(4.2), Inches(1.6), Inches(0.30), TEAL)
add_tb(s2, Inches(7.0), Inches(4.21), Inches(1.4), Inches(0.25),
    "Hyperlink", sz=12, bold=True, color=WHITE)

link_box = add_rect(s2, Inches(6.9), Inches(4.6), Inches(6.0), Inches(0.8), LIGHT_BG)
link_box.line.color.rgb = TEAL
link_box.line.width = Pt(1.5)

demo_html = os.path.join(base_dir, "demo_app.html")
add_tb(s2, Inches(7.1), Inches(4.7), Inches(5.6), Inches(0.5),
    f"Live Demo: {demo_html}\n(Open in Edge/Chrome to interact)",
    sz=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Planned Features
add_rect(s2, Inches(6.9), Inches(5.6), Inches(2.0), Inches(0.30), TEAL)
add_tb(s2, Inches(7.0), Inches(5.61), Inches(1.8), Inches(0.25),
    "Planned Features", sz=12, bold=True, color=WHITE)

add_bullets(s2, Inches(6.9), Inches(5.95), Inches(6.0), Inches(0.9), [
    "Interactive map-based store discovery (GPS-enabled) with real-time inventory",
    "Search & filter by style, size, category, price range; payment gateway integration",
    "Merchant dashboard for inventory mgmt, order fulfillment & sales analytics",
    "Push notifications for order status updates and promotions",
], sz=12)

make_footer(s2)

# ================================================================
# SLIDE 3 — Proof on the field
# ================================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3)
make_header(s3, "Slide 3 of 3")

add_tb(s3, Inches(0.4), Inches(0.6), Inches(12.5), Inches(0.6),
    "Proof on the field", sz=28, bold=True, color=TEAL)

# Status badge
badge = add_rect(s3, Inches(4.0), Inches(1.2), Inches(5.3), Inches(0.55), ORANGE)
add_tb(s3, Inches(4.0), Inches(1.23), Inches(5.3), Inches(0.5),
    "STATUS:  Functional Prototype / Pre-Deployment", sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# --- LEFT: Testimonial ---
add_rect(s3, Inches(0.4), Inches(2.0), Inches(3.0), Inches(0.30), TEAL)
add_tb(s3, Inches(0.5), Inches(2.01), Inches(2.8), Inches(0.25),
    "Customer / User Testimonial", sz=12, bold=True, color=WHITE)

test_box = add_rect(s3, Inches(0.4), Inches(2.4), Inches(5.8), Inches(2.0), LIGHT_BG)
test_box.line.color.rgb = TEAL
test_box.line.width = Pt(1.5)

add_tb(s3, Inches(0.7), Inches(2.6), Inches(5.2), Inches(1.2),
    "\u201cThe prototype gives a clear sense of how the app will "
    "work in real life \u2014 browsing stores nearby, checking stock, "
    "and ordering all in one place. This solves a real pain point "
    "for local shopping.\u201d",
    sz=12, color=DARK)
add_tb(s3, Inches(0.7), Inches(3.8), Inches(5.2), Inches(0.3),
    "\u2014 Early tester feedback (pilot onboarding Q3 2026)",
    sz=10, color=GRAY)

# --- RIGHT: Prototype in use ---
add_rect(s3, Inches(6.8), Inches(2.0), Inches(3.5), Inches(0.30), TEAL)
add_tb(s3, Inches(6.9), Inches(2.01), Inches(3.3), Inches(0.25),
    "Prototype Demo Screens", sz=12, bold=True, color=WHITE)

field_box = add_rect(s3, Inches(6.8), Inches(2.4), Inches(6.1), Inches(2.0), LIGHT_BG)
field_box.line.color.rgb = TEAL
field_box.line.width = Pt(1.5)

# Show multiple small screenshots
add_img(s3, os.path.join(base_dir, "screenshot_home.png"), Inches(7.0), Inches(2.55), Inches(1.5), Inches(1.7))
add_img(s3, os.path.join(base_dir, "screenshot_store.png"), Inches(8.7), Inches(2.55), Inches(1.5), Inches(1.7))
add_img(s3, os.path.join(base_dir, "screenshot_cart.png"), Inches(10.4), Inches(2.55), Inches(1.5), Inches(1.7))

add_tb(s3, Inches(6.8), Inches(4.5), Inches(6.1), Inches(0.3),
    "Interactive prototype showcasing all key user flows",
    sz=10, color=GRAY, align=PP_ALIGN.CENTER)

# --- Next Steps ---
add_rect(s3, Inches(0.4), Inches(4.7), Inches(1.6), Inches(0.30), TEAL)
add_tb(s3, Inches(0.5), Inches(4.71), Inches(1.4), Inches(0.25),
    "Next Steps", sz=12, bold=True, color=WHITE)

add_bullets(s3, Inches(0.4), Inches(5.1), Inches(12.5), Inches(1.7), [
    "1. Functional prototype completed \u2014 demo_app.html showcases full UX flow (discovery, inventory, cart, merchant dashboard)",
    "2. Next: Convert prototype to actual mobile app (React Native / Flutter) with real backend API and database",
    "3. Onboard 10\u201315 pilot stores in a target locality; gather user feedback and iterate on product features",
    "4. Scale to 50+ stores and 1,000+ users in Year 1; capture field photos, testimonials, and usage metrics",
], sz=12, sp=Pt(6))

make_footer(s3)

# Save
output_path = os.path.join(base_dir, "Local_Clothing_Store_Discovery.pptx")
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Size: {os.path.getsize(output_path)} bytes")
